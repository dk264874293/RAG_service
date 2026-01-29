"""
工业级文档处理流水线
支持：PDF、Word、Excel、PPT、HTML、Markdown、图片OCR等
"""

import asyncio
import hashlib
import logging
import aiofiles
from typing import List, Dict, Optional
from pathlib import Path
from ..models.document import Document
from ..extractor.pdf_extractor import PdfExtractor
from ..extractor.word_extractor import WordExtractor
from .adaptive_chunker import AdaptiveChunker

logger = logging.getLogger(__name__)


class UnsupportedFormatError(Exception):
    """不支持的文件格式异常"""

    pass


class DocumentProcessingPipeline:
    """文档处理流水线"""

    def __init__(self, config: Optional[Dict] = None):
        self.config = config or {}
        self.chunker = AdaptiveChunker(config)
        self.supported_formats = {
            ".pdf": self._process_pdf,
            ".docx": self._process_docx,
            ".doc": self._process_doc,
            ".txt": self._process_text,
            ".md": self._process_markdown,
            ".html": self._process_html,
            ".pptx": self._process_pptx,
            ".xlsx": self._process_xlsx,
        }

    def _detect_format(self, file_path: str) -> str:
        """检测文件格式"""
        return Path(file_path).suffix.lower()

    async def process_document(
        self, file_path: str, metadata: Optional[Dict] = None
    ) -> List[Document]:
        """异步处理文档"""
        # 1. 格式检测与验证
        file_ext = self._detect_format(file_path)
        if file_ext not in self.supported_formats:
            raise UnsupportedFormatError(f"不支持的文件格式: {file_ext}")

        # 2. 内容提取
        processor_func = self.supported_formats[file_ext]
        content = await processor_func(file_path)

        # 3. 文本清洗和标准化
        cleaned_content = self._clean_content(content)

        # 4. 分块策略
        enable_chunking = self.config.get("enable_chunking", False)
        if enable_chunking:
            doc_type = self._detect_doc_type(file_path)
            chunks = self.chunker.chunk_document(cleaned_content, doc_type)
        else:
            # 不启用分块，使用完整内容
            chunks = [cleaned_content]

        # 5. 元数据增强
        enhanced_metadata = await self._enhance_metadata(
            metadata or {}, file_path, file_ext
        )

        # 创建文档对象列表
        documents = [
            Document(page_content=chunk, metadata=enhanced_metadata) for chunk in chunks
        ]

        logger.info(f"文档处理完成: 共 {len(documents)} 个分块")
        return documents

    def _detect_doc_type(self, file_path: str) -> str:
        """根据文件路径检测文档类型"""
        from pathlib import Path

        filename = Path(file_path).name.lower()

        if "论文" in filename or "paper" in filename:
            return "research_paper"
        elif "合同" in filename or "contract" in filename or "法律" in filename:
            return "legal_document"
        elif "技术" in filename or "manual" in filename or "api" in filename:
            return "technical_doc"
        elif "财务" in filename or "financial" in filename or "报表" in filename:
            return "financial_report"
        elif ".py" in filename or ".js" in filename or ".java" in filename:
            return "source_code"
        else:
            return "default"

    async def _process_pdf(self, file_path: str) -> str:
        """处理PDF文件"""
        # 根据配置选择是否启用OCR
        enable_ocr = self.config.get("enable_pdf_ocr", True)

        # 统一使用增强版PDF提取器
        extractor = PdfExtractor(
            file_path,
            tenant_id="default",
            user_id="default",
            config=self.config,
            enable_ocr=enable_ocr,
        )
        logger.info(f"使用增强版PDF提取器 (OCR: {enable_ocr}")

        # 在单独的线程中执行同步提取操作
        documents = await asyncio.to_thread(extractor.extract)

        # 记录OCR统计信息（如果可用）
        if enable_ocr and hasattr(extractor, "get_ocr_stats"):
            stats = extractor.get_ocr_stats()
            logger.info(f"PDF处理完成，OCR统计: {stats}")

        # 合并所有文档内容
        return "\n\n".join([doc.page_content for doc in documents])

    async def _process_docx(self, file_path: str) -> str:
        """处理Word文档"""
        extractor = WordExtractor(file_path, tenant_id="default", user_id="default")
        # 在单独的线程中执行同步提取操作
        documents = await asyncio.to_thread(extractor.extract)
        return "\n\n".join([doc.page_content for doc in documents])

    async def _process_doc(self, file_path: str) -> str:
        """处理旧版Word文档"""
        # 可以调用WordExtractor或使用其他库
        return await self._process_docx(file_path)

    async def _process_text(self, file_path: str) -> str:
        """处理纯文本文件"""
        async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
            return await f.read()

    async def _process_markdown(self, file_path: str) -> str:
        """处理Markdown文件"""
        async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
            return await f.read()

    async def _process_html(self, file_path: str) -> str:
        """处理HTML文件"""

        # 定义同步处理函数
        def _parse_html_sync(content: str) -> str:
            try:
                from bs4 import BeautifulSoup

                soup = BeautifulSoup(content, "html.parser")
                return soup.get_text()
            except ImportError:
                # 如果没有BeautifulSoup，简单提取文本
                import re

                return re.sub(r"<[^>]+>", "", content)

        # 异步读取文件内容
        async with aiofiles.open(file_path, "r", encoding="utf-8") as f:
            content = await f.read()

        # 在单独的线程中执行HTML解析（CPU密集型操作）
        return await asyncio.to_thread(_parse_html_sync, content)

    async def _process_pptx(self, file_path: str) -> str:
        """处理PowerPoint文件"""

        # 定义同步处理函数
        def _parse_pptx_sync(file_path: str) -> str:
            try:
                from pptx import Presentation

                prs = Presentation(file_path)
                text_content = []
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text_content.append(shape.text)
                return "\n\n".join(text_content)
            except ImportError:
                raise UnsupportedFormatError("需要安装python-pptx库来处理PPTX文件")

        # 在单独的线程中执行PPTX解析（CPU密集型操作）
        return await asyncio.to_thread(_parse_pptx_sync, file_path)

    async def _process_xlsx(self, file_path: str) -> str:
        """处理Excel文件"""

        # 定义同步处理函数
        def _parse_xlsx_sync(file_path: str) -> str:
            try:
                import pandas as pd

                df = pd.read_excel(file_path, sheet_name=None)
                text_content = []
                for sheet_name, sheet_df in df.items():
                    text_content.append(f"工作表: {sheet_name}\n{sheet_df.to_string()}")
                return "\n\n".join(text_content)
            except ImportError:
                raise UnsupportedFormatError(
                    "需要安装pandas和openpyxl库来处理Excel文件"
                )

        # 在单独的线程中执行Excel解析（CPU密集型操作）
        return await asyncio.to_thread(_parse_xlsx_sync, file_path)

    def _clean_content(self, content: str) -> str:
        """文本清洗和标准化"""
        # 去除多余的空白字符
        import re

        content = re.sub(r"\s+", " ", content)
        # 去除特殊字符（可根据需要调整）
        content = content.strip()
        return content

    async def _enhance_metadata(
        self, metadata: Dict, file_path: str, file_ext: str
    ) -> Dict:
        """元数据增强"""
        enhanced = metadata.copy()
        path_obj = Path(file_path)

        # 添加文件信息
        file_hash = None
        if path_obj.exists():
            file_hash = await self._calculate_file_hash(file_path)

        enhanced.update(
            {
                "source": str(file_path),
                "filename": path_obj.name,
                "file_type": file_ext,
                "file_size": path_obj.stat().st_size if path_obj.exists() else 0,
                "file_hash": file_hash,
            }
        )

        return enhanced

    async def _calculate_file_hash(self, file_path: str) -> str:
        """异步计算文件哈希值"""
        hash_md5 = hashlib.md5()
        async with aiofiles.open(file_path, "rb") as f:
            while True:
                chunk = await f.read(4096)
                if not chunk:
                    break
                hash_md5.update(chunk)
        return hash_md5.hexdigest()
